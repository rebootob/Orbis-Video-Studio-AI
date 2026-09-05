import io
import uuid
import pytest
import fitz
import docx
import pptx
from app.models.project import Project
from app.models.asset import Asset
from app.models.document_extraction import DocumentExtraction
from app.core.config import settings


def create_sample_pdf_bytes(text: str = "Sample PDF Text") -> bytes:
    doc = fitz.open()
    page = doc.new_page()
    page.insert_text((50, 50), text)
    pdf_bytes = doc.tobytes()
    doc.close()
    return pdf_bytes


def create_empty_pdf_bytes() -> bytes:
    doc = fitz.open()
    doc.new_page()  # Blank page without text layer
    pdf_bytes = doc.tobytes()
    doc.close()
    return pdf_bytes


def create_sample_docx_bytes(paragraphs: list) -> bytes:
    doc = docx.Document()
    for p in paragraphs:
        doc.add_paragraph(p)
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def create_sample_pptx_bytes(slide_texts: list) -> bytes:
    prs = pptx.Presentation()
    blank_slide_layout = prs.slide_layouts[6]
    for text in slide_texts:
        slide = prs.slides.add_slide(blank_slide_layout)
        tx_box = slide.shapes.add_textbox(0, 0, 100, 100)
        tf = tx_box.text_frame
        tf.text = text
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def test_txt_and_markdown_extraction_with_unicode(client, db_session):
    project = Project(title="Unicode Doc Project", status="DRAFT")
    db_session.add(project)
    db_session.commit()

    thai_english_text = (
        "สวัสดีโลก Hello World!\n\n"
        "นี่คือการทดสอบระบบสกัดข้อความ Orbis Video Studio AI.\n"
        "Japanese test: こんにちは世界\n"
        "Mixed script paragraph: Cinematic Storyboard Brief 2026."
    )

    # 1. Upload TXT Asset
    upload_resp = client.post(
        "/api/v1/assets/upload",
        data={"project_id": str(project.id), "name": "Thai Brief", "asset_type": "DOCUMENT"},
        files={"file": ("brief_th.txt", io.BytesIO(thai_english_text.encode("utf-8")), "text/plain")},
    )
    assert upload_resp.status_code == 201
    asset_id = upload_resp.json()["id"]

    # 2. Extract Document
    ext_resp = client.post(f"/api/v1/assets/{asset_id}/extract")
    assert ext_resp.status_code == 200
    data = ext_resp.json()

    assert data["asset_id"] == asset_id
    assert data["document_type"] == "txt"
    assert data["status"] == "SUCCESS"
    assert "สวัสดีโลก Hello World!" in data["extracted_text"]
    assert "こんにちは世界" in data["extracted_text"]
    assert data["character_count"] > 0
    assert data["segment_count"] >= 2
    assert data["extraction_duration_ms"] >= 0.0
    assert data["extraction_method"] == "text-decoder"

    # 3. GET extraction endpoint
    get_ext = client.get(f"/api/v1/assets/{asset_id}/extraction")
    assert get_ext.status_code == 200
    assert get_ext.json()["id"] == data["id"]


def test_docx_extraction(client, db_session):
    project = Project(title="DOCX Project", status="DRAFT")
    db_session.add(project)
    db_session.commit()

    docx_bytes = create_sample_docx_bytes([
        "Chapter 1: The Space Station",
        "The station orbited silently above Jupiter.",
        "Captain Sara adjusted the navigation controls."
    ])

    upload_resp = client.post(
        "/api/v1/assets/upload",
        data={"project_id": str(project.id), "name": "Screenplay Doc", "asset_type": "DOCUMENT"},
        files={"file": ("screenplay.docx", io.BytesIO(docx_bytes), "application/vnd.openxmlformats-officedocument.wordprocessingml.document")},
    )
    asset_id = upload_resp.json()["id"]

    ext_resp = client.post(f"/api/v1/assets/{asset_id}/extract")
    assert ext_resp.status_code == 200
    data = ext_resp.json()

    assert data["document_type"] == "docx"
    assert data["status"] == "SUCCESS"
    assert "Chapter 1: The Space Station" in data["extracted_text"]
    assert "Sara adjusted the navigation controls" in data["extracted_text"]
    assert data["segment_count"] == 3
    assert data["extraction_method"] == "python-docx"


def test_pptx_extraction(client, db_session):
    project = Project(title="PPTX Project", status="DRAFT")
    db_session.add(project)
    db_session.commit()

    pptx_bytes = create_sample_pptx_bytes([
        "Orbis Pitch Presentation Slide 1",
        "Key Storyboard Features & Pipeline Slide 2"
    ])

    upload_resp = client.post(
        "/api/v1/assets/upload",
        data={"project_id": str(project.id), "name": "Pitch Deck", "asset_type": "DOCUMENT"},
        files={"file": ("pitch.pptx", io.BytesIO(pptx_bytes), "application/vnd.openxmlformats-officedocument.presentationml.presentation")},
    )
    asset_id = upload_resp.json()["id"]

    ext_resp = client.post(f"/api/v1/assets/{asset_id}/extract")
    assert ext_resp.status_code == 200
    data = ext_resp.json()

    assert data["document_type"] == "pptx"
    assert data["status"] == "SUCCESS"
    assert "Slide 1" in data["extracted_text"]
    assert "Slide 2" in data["extracted_text"]
    assert data["segment_count"] == 2
    assert data["extraction_method"] == "python-pptx"


def test_pdf_text_layer_extraction(client, db_session):
    project = Project(title="PDF Project", status="DRAFT")
    db_session.add(project)
    db_session.commit()

    pdf_bytes = create_sample_pdf_bytes("PDF Scene Heading: EXT. MOON BASE - NIGHT")

    upload_resp = client.post(
        "/api/v1/assets/upload",
        data={"project_id": str(project.id), "name": "Moon Base Script", "asset_type": "DOCUMENT"},
        files={"file": ("moon_base.pdf", io.BytesIO(pdf_bytes), "application/pdf")},
    )
    asset_id = upload_resp.json()["id"]

    ext_resp = client.post(f"/api/v1/assets/{asset_id}/extract")
    assert ext_resp.status_code == 200
    data = ext_resp.json()

    assert data["document_type"] == "pdf"
    assert data["status"] == "SUCCESS"
    assert "EXT. MOON BASE - NIGHT" in data["extracted_text"]
    assert data["extraction_method"] == "PyMuPDF"


def test_pdf_no_text_layer_returns_warning(client, db_session):
    project = Project(title="Empty PDF Project", status="DRAFT")
    db_session.add(project)
    db_session.commit()

    empty_pdf_bytes = create_empty_pdf_bytes()

    upload_resp = client.post(
        "/api/v1/assets/upload",
        data={"project_id": str(project.id)},
        files={"file": ("scanned.pdf", io.BytesIO(empty_pdf_bytes), "application/pdf")},
    )
    asset_id = upload_resp.json()["id"]

    ext_resp = client.post(f"/api/v1/assets/{asset_id}/extract")
    assert ext_resp.status_code == 200
    data = ext_resp.json()

    assert data["status"] == "NO_TEXT_LAYER"
    assert any("OCR_REQUIRED" in w for w in data["warnings"])


def test_unsupported_document_type(client, db_session):
    project = Project(title="Binary Project", status="DRAFT")
    db_session.add(project)
    db_session.commit()

    upload_resp = client.post(
        "/api/v1/assets/upload",
        data={"project_id": str(project.id)},
        files={"file": ("executable.exe", io.BytesIO(b"MZ12345"), "application/octet-stream")},
    )
    asset_id = upload_resp.json()["id"]

    ext_resp = client.post(f"/api/v1/assets/{asset_id}/extract")
    assert ext_resp.status_code == 400
    assert "unsupported document format" in ext_resp.json()["detail"].lower()


def test_missing_asset_and_extraction(client):
    fake_asset_id = str(uuid.uuid4())

    post_resp = client.post(f"/api/v1/assets/{fake_asset_id}/extract")
    assert post_resp.status_code == 404

    get_resp = client.get(f"/api/v1/assets/{fake_asset_id}/extraction")
    assert get_resp.status_code == 404


def test_re_extraction_force_flag(client, db_session):
    project = Project(title="Force Extract Project", status="DRAFT")
    db_session.add(project)
    db_session.commit()

    upload_resp = client.post(
        "/api/v1/assets/upload",
        data={"project_id": str(project.id)},
        files={"file": ("doc.txt", io.BytesIO(b"Initial text payload"), "text/plain")},
    )
    asset_id = upload_resp.json()["id"]

    # First extraction
    ext1 = client.post(f"/api/v1/assets/{asset_id}/extract").json()

    # Second extraction without force (returns same cached record)
    ext2 = client.post(f"/api/v1/assets/{asset_id}/extract").json()
    assert ext1["id"] == ext2["id"]

    # Third extraction with force=True
    ext3 = client.post(f"/api/v1/assets/{asset_id}/extract?force=true").json()
    assert ext3["id"] == ext1["id"]


def test_document_size_limit_validation(client, db_session, monkeypatch):
    project = Project(title="Large Doc Project", status="DRAFT")
    db_session.add(project)
    db_session.commit()

    # Temporarily set max document bytes to 50 bytes for test
    monkeypatch.setattr(settings, "MAX_DOCUMENT_BYTES", 50)

    upload_resp = client.post(
        "/api/v1/assets/upload",
        data={"project_id": str(project.id)},
        files={"file": ("huge_file.txt", io.BytesIO(b"A" * 100), "text/plain")},
    )
    asset_id = upload_resp.json()["id"]

    ext_resp = client.post(f"/api/v1/assets/{asset_id}/extract")
    assert ext_resp.status_code == 400
    assert "exceeds maximum limit" in ext_resp.json()["detail"].lower()


def test_pre_download_size_guard_prevents_storage_get_object(client, db_session, monkeypatch):
    from unittest.mock import MagicMock
    from app.services.document_extraction.service import DocumentExtractionService, DocumentExtractionError

    project = Project(title="Pre-download Guard Project", status="DRAFT")
    db_session.add(project)
    db_session.commit()

    asset = Asset(
        id=uuid.uuid4(),
        project_id=project.id,
        name="Pre-download Guard Asset",
        asset_type="DOCUMENT",
        original_filename="large.txt",
        content_type="text/plain",
        file_size_bytes=settings.MAX_DOCUMENT_BYTES + 1000,
        checksum_sha256="dummychecksum1234567890",
        storage_bucket="test-bucket",
        storage_key="test-key",
    )
    db_session.add(asset)
    db_session.commit()

    mock_storage = MagicMock()
    service = DocumentExtractionService(db=db_session, storage=mock_storage)

    with pytest.raises(DocumentExtractionError) as exc_info:
        service.extract_asset_document(asset.id)

    assert exc_info.value.code == "DOCUMENT_TOO_LARGE"
    mock_storage.get_object.assert_not_called()


def test_pdf_hard_character_limit(client, db_session, monkeypatch):
    monkeypatch.setattr(settings, "MAX_EXTRACTED_CHARACTERS", 30)

    project = Project(title="PDF Limit Project", status="DRAFT")
    db_session.add(project)
    db_session.commit()

    pdf_bytes = create_sample_pdf_bytes("This is a long PDF text that exceeds 30 characters easily.")

    upload_resp = client.post(
        "/api/v1/assets/upload",
        data={"project_id": str(project.id)},
        files={"file": ("limit.pdf", io.BytesIO(pdf_bytes), "application/pdf")},
    )
    asset_id = upload_resp.json()["id"]

    ext_resp = client.post(f"/api/v1/assets/{asset_id}/extract")
    assert ext_resp.status_code == 200
    data = ext_resp.json()

    assert data["character_count"] <= 30
    assert len(data["extracted_text"]) <= 30


def test_docx_hard_character_limit_and_table_loop_termination(client, db_session, monkeypatch):
    monkeypatch.setattr(settings, "MAX_EXTRACTED_CHARACTERS", 25)

    project = Project(title="DOCX Limit Project", status="DRAFT")
    db_session.add(project)
    db_session.commit()

    doc = docx.Document()
    doc.add_paragraph("Paragraph 1 with long text exceeding limit.")
    t1 = doc.add_table(rows=2, cols=2)
    t1.rows[0].cells[0].text = "Cell 1 text"
    t1.rows[0].cells[1].text = "Cell 2 text"
    t2 = doc.add_table(rows=2, cols=2)
    t2.rows[0].cells[0].text = "Cell 3 text"
    t2.rows[0].cells[1].text = "Cell 4 text"

    buf = io.BytesIO()
    doc.save(buf)
    docx_bytes = buf.getvalue()

    upload_resp = client.post(
        "/api/v1/assets/upload",
        data={"project_id": str(project.id)},
        files={"file": ("table_limit.docx", io.BytesIO(docx_bytes), "application/vnd.openxmlformats-officedocument.wordprocessingml.document")},
    )
    asset_id = upload_resp.json()["id"]

    ext_resp = client.post(f"/api/v1/assets/{asset_id}/extract")
    assert ext_resp.status_code == 200
    data = ext_resp.json()

    assert data["character_count"] <= 25
    assert len(data["extracted_text"]) <= 25


def test_pptx_hard_character_limit(client, db_session, monkeypatch):
    monkeypatch.setattr(settings, "MAX_EXTRACTED_CHARACTERS", 20)

    project = Project(title="PPTX Limit Project", status="DRAFT")
    db_session.add(project)
    db_session.commit()

    pptx_bytes = create_sample_pptx_bytes([
        "Slide 1 long text content exceeding max limit",
        "Slide 2 text content"
    ])

    upload_resp = client.post(
        "/api/v1/assets/upload",
        data={"project_id": str(project.id)},
        files={"file": ("limit.pptx", io.BytesIO(pptx_bytes), "application/vnd.openxmlformats-officedocument.presentationml.presentation")},
    )
    asset_id = upload_resp.json()["id"]

    ext_resp = client.post(f"/api/v1/assets/{asset_id}/extract")
    assert ext_resp.status_code == 200
    data = ext_resp.json()

    assert data["character_count"] <= 20
    assert len(data["extracted_text"]) <= 20


def test_storage_access_failed_error_semantics(client, db_session, mock_storage):
    from unittest.mock import MagicMock

    project = Project(title="Storage Failure Project", status="DRAFT")
    db_session.add(project)
    db_session.commit()

    asset = Asset(
        id=uuid.uuid4(),
        project_id=project.id,
        name="Failed Storage Asset",
        asset_type="DOCUMENT",
        original_filename="test.txt",
        content_type="text/plain",
        file_size_bytes=100,
        checksum_sha256="dummychecksum1234567890",
        storage_bucket="test-bucket",
        storage_key="test-key",
    )
    db_session.add(asset)
    db_session.commit()

    mock_storage.get_object = MagicMock(side_effect=RuntimeError("S3 Connection Refused / Provider Failure"))

    ext_resp = client.post(f"/api/v1/assets/{asset.id}/extract")
    assert ext_resp.status_code == 500
    assert "Failed to access object storage payload" in ext_resp.json()["detail"]

