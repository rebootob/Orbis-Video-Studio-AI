import io
import time
import uuid
from typing import Dict, Any, List
import pptx
from app.core.config import settings
from app.services.document_extraction.base import DocumentExtractor, ExtractionResult


class PptxDocumentExtractor(DocumentExtractor):
    """Fast native PowerPoint presentation (.pptx) extractor using python-pptx."""

    def supports(self, document_type: str) -> bool:
        return document_type == "pptx"

    def extract(self, asset_id: uuid.UUID, content: bytes, metadata: Dict[str, Any]) -> ExtractionResult:
        start_time = time.perf_counter()
        warnings: List[str] = []

        try:
            prs = pptx.Presentation(io.BytesIO(content))
        except Exception as e:
            duration_ms = round((time.perf_counter() - start_time) * 1000, 2)
            return ExtractionResult(
                asset_id=asset_id,
                document_type="pptx",
                extracted_text="",
                character_count=0,
                segment_count=0,
                extraction_status="FAILED",
                extraction_method="python-pptx",
                extraction_duration_ms=duration_ms,
                error_message=f"Failed to parse PPTX presentation: {str(e)}",
            )

        total_slides = len(prs.slides)
        if total_slides > settings.MAX_DOCUMENT_PAGES:
            warnings.append(
                f"Presentation slide count ({total_slides}) exceeds limit ({settings.MAX_DOCUMENT_PAGES}). Extracted first {settings.MAX_DOCUMENT_PAGES} slides."
            )
            slides_to_read = settings.MAX_DOCUMENT_PAGES
        else:
            slides_to_read = total_slides

        slide_texts: List[str] = []
        segments: List[Dict[str, Any]] = []
        total_chars = 0

        for slide_idx in range(slides_to_read):
            slide = prs.slides[slide_idx]
            shape_texts: List[str] = []

            for shape in slide.shapes:
                if shape.has_text_frame:
                    tf_text = shape.text_frame.text.strip()
                    if tf_text:
                        clean_tf = tf_text.replace("\r\n", "\n").replace("\r", "\n").replace("\x00", "")
                        shape_texts.append(clean_tf)

            slide_text = "\n".join(shape_texts).strip()
            if slide_text:
                slide_texts.append(f"--- Slide {slide_idx + 1} ---\n{slide_text}")
                segments.append({
                    "index": slide_idx + 1,
                    "type": "slide",
                    "text": slide_text,
                })
                total_chars += len(slide_text)

                if total_chars >= settings.MAX_EXTRACTED_CHARACTERS:
                    warnings.append(f"Character extraction limit ({settings.MAX_EXTRACTED_CHARACTERS}) reached.")
                    break

        combined_text = "\n\n".join(slide_texts)
        duration_ms = round((time.perf_counter() - start_time) * 1000, 2)
        status = "SUCCESS" if combined_text.strip() else "NO_TEXT_LAYER"

        return ExtractionResult(
            asset_id=asset_id,
            document_type="pptx",
            extracted_text=combined_text,
            character_count=len(combined_text),
            segment_count=len(segments),
            segments=segments,
            extraction_status=status,
            extraction_method="python-pptx",
            extraction_duration_ms=duration_ms,
            warnings=warnings,
        )
